import os
import sys
import subprocess
import re
import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
MD_REPORT_PATH = os.path.join(DOCS_DIR, "Microservices_Technical_Report_30_50_Pages.md")
HTML_REPORT_PATH = os.path.join(DOCS_DIR, "Microservices_Technical_Report_30_50_Pages.html")
PDF_REPORT_PATH = os.path.join(DOCS_DIR, "Microservices_Technical_Report_30_50_Pages.pdf")
SLIDES_MD_PATH = os.path.join(DOCS_DIR, "EasyService_Presentation_20_Slides.md")
SLIDES_PPTX_PATH = os.path.join(DOCS_DIR, "EasyService_Presentation_20_Slides.pptx")

print("Starting EasyService PDF & 20-Slide PPTX Generator...")

# ---------------------------------------------------------
# STEP 1: CONVERT TECHNICAL REPORT TO COLORFUL HTML & PDF
# ---------------------------------------------------------

with open(MD_REPORT_PATH, "r", encoding="utf-8") as f:
    md_content = f.read()

# Convert markdown to html using markdown library with table and code extensions
html_body = markdown.markdown(md_content, extensions=['tables', 'fenced_code', 'toc'])

css_styles = """
<style>
  @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;600;700;800&family=Fira+Code:wght@400;600&display=swap');

  @page {
    size: A4 portrait;
    margin: 20mm 15mm 20mm 15mm;
  }

  body {
    font-family: 'Plus Jakarta Sans', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
    color: #1e293b;
    background-color: #ffffff;
    line-height: 1.6;
    font-size: 10.5pt;
  }

  /* Cover Page Styling */
  .cover-container {
    background: linear-gradient(135deg, #0f172a 0%, #1e293b 50%, #1e1b4b 100%);
    color: #ffffff;
    padding: 50px 40px;
    border-radius: 16px;
    margin-bottom: 40px;
    page-break-after: always;
    box-shadow: 0 10px 25px rgba(0,0,0,0.2);
  }

  .cover-badge {
    display: inline-block;
    background: linear-gradient(90deg, #d4af37, #f59e0b);
    color: #0f172a;
    font-size: 10pt;
    font-weight: 800;
    padding: 6px 16px;
    border-radius: 20px;
    text-transform: uppercase;
    letter-spacing: 1.5px;
    margin-bottom: 20px;
  }

  .cover-title {
    font-size: 26pt;
    font-weight: 800;
    color: #ffffff;
    line-height: 1.2;
    margin-bottom: 15px;
    border-bottom: 3px solid #d4af37;
    padding-bottom: 15px;
  }

  .cover-subtitle {
    font-size: 13pt;
    color: #94a3b8;
    margin-bottom: 35px;
    font-weight: 500;
  }

  .meta-grid {
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 20px;
    background: rgba(255, 255, 255, 0.05);
    padding: 20px;
    border-radius: 12px;
    border: 1px solid rgba(255, 255, 255, 0.1);
    margin-bottom: 30px;
  }

  .meta-item strong {
    color: #d4af37;
    display: block;
    font-size: 9pt;
    text-transform: uppercase;
    letter-spacing: 1px;
  }

  .meta-item span {
    color: #f8fafc;
    font-size: 11pt;
    font-weight: 600;
  }

  /* Member Table in Cover */
  .member-table {
    width: 100%;
    border-collapse: collapse;
    margin-top: 15px;
    background: rgba(255,255,255,0.03);
    border-radius: 8px;
    overflow: hidden;
  }

  .member-table th {
    background: #d4af37;
    color: #0f172a;
    font-size: 9pt;
    font-weight: 800;
    text-transform: uppercase;
    padding: 10px 14px;
    text-align: left;
  }

  .member-table td {
    padding: 9px 14px;
    border-bottom: 1px solid rgba(255,255,255,0.08);
    color: #e2e8f0;
    font-size: 9.5pt;
  }

  /* Document Typography */
  h1 {
    font-size: 18pt;
    font-weight: 800;
    color: #0f172a;
    border-left: 5px solid #d4af37;
    padding-left: 12px;
    margin-top: 35px;
    margin-bottom: 15px;
    page-break-before: always;
  }

  h2 {
    font-size: 14pt;
    font-weight: 700;
    color: #1e293b;
    border-bottom: 2px solid #e2e8f0;
    padding-bottom: 6px;
    margin-top: 25px;
    margin-bottom: 12px;
  }

  h3 {
    font-size: 11.5pt;
    font-weight: 700;
    color: #334155;
    margin-top: 18px;
    margin-bottom: 8px;
  }

  p {
    margin-bottom: 12px;
    text-align: justify;
  }

  ul, ol {
    margin-bottom: 14px;
    padding-left: 24px;
  }

  li {
    margin-bottom: 4px;
  }

  /* Tables Styling */
  table {
    width: 100%;
    border-collapse: collapse;
    margin: 18px 0;
    font-size: 9.5pt;
    page-break-inside: avoid;
  }

  th {
    background-color: #0f172a;
    color: #ffffff;
    font-weight: 700;
    text-align: left;
    padding: 10px 12px;
    border: 1px solid #0f172a;
  }

  td {
    padding: 8px 12px;
    border: 1px solid #cbd5e1;
    color: #334155;
  }

  tr:nth-child(even) {
    background-color: #f8fafc;
  }

  /* Code Block Styling */
  pre {
    background-color: #0f172a;
    color: #f8fafc;
    padding: 14px 18px;
    border-radius: 8px;
    font-family: 'Fira Code', Consolas, monospace;
    font-size: 8.5pt;
    line-height: 1.5;
    overflow-x: auto;
    border-left: 4px solid #d4af37;
    margin: 16px 0;
    page-break-inside: avoid;
  }

  code {
    font-family: 'Fira Code', Consolas, monospace;
    background-color: #f1f5f9;
    color: #b45309;
    padding: 2px 6px;
    border-radius: 4px;
    font-size: 9pt;
  }

  pre code {
    background-color: transparent;
    color: inherit;
    padding: 0;
  }

  /* Horizontal Rule */
  hr {
    border: none;
    height: 1px;
    background: #e2e8f0;
    margin: 25px 0;
  }

  /* Callout box */
  blockquote {
    background: #fffbebfb;
    border-left: 4px solid #f59e0b;
    margin: 16px 0;
    padding: 12px 18px;
    border-radius: 0 8px 8px 0;
    color: #92400e;
    font-size: 10pt;
  }
</style>
"""

full_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <title>EasyService Microservices Architecture Technical Report</title>
  {css_styles}
</head>
<body>
  {html_body}
</body>
</html>
"""

with open(HTML_REPORT_PATH, "w", encoding="utf-8") as f:
    f.write(full_html)

print(f"[OK] Generated HTML Report at: {HTML_REPORT_PATH}")

# Convert HTML to PDF using Microsoft Edge Headless
edge_path = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
if os.path.exists(edge_path):
    cmd = [
        edge_path,
        "--headless",
        "--disable-gpu",
        f"--print-to-pdf={PDF_REPORT_PATH}",
        HTML_REPORT_PATH
    ]
    subprocess.run(cmd, check=True)
    print(f"[OK] Generated Colorful PDF Report at: {PDF_REPORT_PATH}")
else:
    print("Warning: Microsoft Edge executable not found at standard path. HTML report generated.")


# ---------------------------------------------------------
# STEP 2: GENERATE 20-SLIDE PRESENTATION IN MARKDOWN & PPTX
# ---------------------------------------------------------

slides_data = [
    {
        "num": 1,
        "title": "EasyService Architecture & Implementation Overview",
        "subtitle": "Cloud-Native Enterprise Marketplace Platform — Group Two",
        "bullets": [
            "Course: Cloud Computing & Security (Microservices Architecture)",
            "Institution: Addis Ababa University — School of IT & Engineering",
            "Lead Contributor: Abel Seleshe (Architecture & Decomposition Lead)",
            "Team Members: Baheran Tesfaye, Nebiyu Yohannes, Wondesen Teshale, Esrom Basazinaw, Yohannes Seyum",
            "Target Implementation: Java 21, Spring Boot 3.2, Svelte Frontend & Docker/K8s"
        ]
    },
    {
        "num": 2,
        "title": "Executive Summary & System Vision",
        "subtitle": "Transforming Monolithic Service Portals into Autonomous Microservices",
        "bullets": [
            "Vision: Modern Ethiopian marketplace unifying Hotels, Mobility, Crafts, and Local Services",
            "Core Problem: Monolithic bottlenecks, shared database contention, and unified blast radiuses",
            "Architectural Solution: 6 Autonomous Bounded Context Microservices",
            "Security Foundation: Stateless OAuth2/JWT incorporating Fayda National Digital ID claims",
            "Key Capabilities: Reactive Wallet, Dynamic Spin-Wheel Engine & 1-Click Hot Deals Booking"
        ]
    },
    {
        "num": 3,
        "title": "EasyService Monolith Decomposition",
        "subtitle": "Decomposing Legacy Systems into 6 Bounded Contexts",
        "bullets": [
            "User-Service (Identity & Fayda/Passport Verification)",
            "Listing-Service (Marketplace Inventory & Multi-Room Hotel Units)",
            "Booking-Service (Atomic Inventory Locking & Booking Stepper State Engine)",
            "Payment-Service (EasyService Wallet Balance & Telebirr/CBE Transfers)",
            "Notification-Service (Real-time Event Alerts & Pre-Booking Reminders)",
            "Wheel-Promo-Service (Daily Spin Cost Limits & Unique Single-Use Promo Generator)"
        ]
    },
    {
        "num": 4,
        "title": "Domain Model: Provider vs. Bookable Units",
        "subtitle": "Hierarchical Property Inventory Model",
        "bullets": [
            "Root Provider/Business (e.g. Kuriftu Resort & Spa, Skylight Hotel, Habesha Kemis)",
            "Listing Metadata: Location, Amenities, Ratings, Photos & Policies",
            "Bookable Units / Options: Standard Room (3,500 ETB), Deluxe Room (5,000 ETB), Family Villa (10,000 ETB)",
            "Car Fleet Units: Toyota RAV4, Hyundai Elantra, Land Cruiser V8 (Daily Rates)",
            "Store Products: Yirgacheffe Coffee, Silk Kemis, Leather Briefcase (Item Units)"
        ]
    },
    {
        "num": 5,
        "title": "Hotels & Accommodation Marketplace Engine",
        "subtitle": "Kuriftu, Skylight & Premium Resort Booking Flow",
        "bullets": [
            "Discovery Page: Filter by Price, Star Rating, Property Type & Amenities",
            "Property Detail Page: Image Gallery, Location Coordinates & Full Specs",
            "Room Selection Matrix: Direct room selection with real-time room availability counts",
            "Date Pickers: Interactive Check-in & Check-out date range selection",
            "Nights Multiplier: Automated subtotal calculation based on stay duration"
        ]
    },
    {
        "num": 6,
        "title": "Car Rental & Mobility Marketplace Engine",
        "subtitle": "Vehicle Fleet Specs & Rental Duration Engine",
        "bullets": [
            "Fleet Inventory: SUV, Sedan, Luxury & Off-Road categories",
            "Vehicle Specs: Seats, Transmission, Fuel Type, Driver Option & Unlimited KM",
            "Rental Date Range: Pickup Date and Return Date selector",
            "Day Multiplier: Automatic rental duration calculation (Daily Rate × Days × Quantity)",
            "Instant Validation: Return date must strictly be after pickup date"
        ]
    },
    {
        "num": 7,
        "title": "Cultural Products, Crafts & Events Portal",
        "subtitle": "Ethiopian Heritage Kemis, Leather Goods & Ticket Booking",
        "bullets": [
            "Authentic Local Stores: Habesha Heritage Kemis, Entoto Spices, Axum Crafts",
            "Item Variants: Specific sizes, silk borders, and material specifications",
            "Concerts & Cultural Events: Great Ethiopian Run, Jan Meda Festival tickets",
            "Preferred Date Picker: Schedule delivery or event attendance date",
            "Direct Provider Settlement: Integrated commission split engine"
        ]
    },
    {
        "num": 8,
        "title": "Customer Identity & Fayda Verification",
        "subtitle": "Zero-Trust Security & National ID Integration",
        "bullets": [
            "AAU Academic Domain SSO: Restricted registration for verified student/faculty credentials",
            "Fayda National Digital ID Integration: Verification badge (FAYDA_VERIFIED)",
            "Passport Identity Status: Level-2 identity trust scores",
            "Stateless JWT Security: RSA-256 signed bearer tokens",
            "Role-Based Access Control: CUSTOMER, PROVIDER, SYSTEM_ADMIN roles"
        ]
    },
    {
        "num": 9,
        "title": "Reactive EasyService Wallet System",
        "subtitle": "Real-time Balance Management & Atomic Deductions",
        "bullets": [
            "Simulated Wallet Store: $currentUser reactive state management",
            "Instant Balance Check: Blocks transaction if balance < total payable",
            "Atomic Deduction: Mutates user balance instantly across UI navbar & passport",
            "Alternative Payment Channels: Integrated Telebirr, CBE Birr & Card options",
            "Instant Refund Engine: Wallet credits returned upon booking cancellation"
        ]
    },
    {
        "num": 10,
        "title": "Spin-the-Wheel Promotion Engine",
        "subtitle": "Tiered Spin Costs & Unique Single-Use Promo Codes",
        "bullets": [
            "Tiered Cost Rules: Spins 1-3 (FREE), Spins 4-9 (50 ETB), Spins 10-15 (100 ETB)",
            "Wallet Protection: Deducts spin cost directly from EasyService wallet balance",
            "Randomized Code Generator: Generates unique code per spin (e.g. WIN10-8492)",
            "Central Code Registry: Registers code in activePromoCodes store to prevent reuse",
            "Clipboard Copy Action: 1-Click copy button with text fallback"
        ]
    },
    {
        "num": 11,
        "title": "Hot Deals & Dynamic Pricing Engine",
        "subtitle": "Pre-applied Promos & Flash Sale Bookings",
        "bullets": [
            "Marketing Banner Integration: Direct 1-click 'Book Deal →' triggers",
            "Pre-applied Promo Deep Linking: Automatically loads promo codes (SUMMER20, ETHIO30)",
            "Dynamic Savings Calculation: Displays percentage discount & ETB amount saved",
            "Service Fee Inclusion: Transparent 5% platform service fee breakdown",
            "Real-time Subtotal Formula: (Unit Price × Quantity × Duration - Discount) + Fee"
        ]
    },
    {
        "num": 12,
        "title": "End-to-End Booking Stepper Architecture",
        "subtitle": "5-Step Transactional Booking Flow",
        "bullets": [
            "Step 1 (Select): Quantity, Dates (Check-in/out), and Promo Code input",
            "Step 2 (Details): Customer Name, Email, Phone & Identity Status verification",
            "Step 3 (Review): Comprehensive line-item summary and cancellation policy",
            "Step 4 (Payment): Payment method selection (Wallet, Telebirr, CBE) & charge execution",
            "Step 5 (Confirm): Green receipt badge, reference ID, schedule dates & wallet update"
        ]
    },
    {
        "num": 13,
        "title": "Date Range Pickers & Schedule Calculation",
        "subtitle": "Native Date Controls & Duration Engine",
        "bullets": [
            "Native HTML5 Date Pickers: Minimum date bounded to current day (getTodayStr)",
            "Date Difference Engine: Math.ceil((end - start) / 86,400,000)",
            "Category-Aware Labels: Check-in/out for Hotels, Pickup/Return for Cars",
            "Date Range Badge: Visual summary badge displaying computed duration",
            "Confirmed Schedule Payload: Persists dates to bookingResult and store history"
        ]
    },
    {
        "num": 14,
        "title": "Microservices Technical Architecture",
        "subtitle": "Spring Boot 3.2, REST & Messaging Mesh",
        "bullets": [
            "Framework: Java 21 LTS with Spring Boot 3.2 Cloud-Native stack",
            "API Gateway: Spring Cloud Gateway for routing, rate limiting & SSL termination",
            "Service Discovery: Netflix Eureka Service Registry",
            "Inter-Service Sync: RESTful HTTP/2 APIs & gRPC binary protocol",
            "Async Messaging: Apache Kafka event streams for notifications and Saga triggers"
        ]
    },
    {
        "num": 15,
        "title": "Database Schema & Data Isolation",
        "subtitle": "Database per Service Pattern Implementation",
        "bullets": [
            "Isolated Storage: Dedicated PostgreSQL schema per microservice boundary",
            "User DB: Customer profiles, credentials & identity verification flags",
            "Listing DB: Provider properties, room units, specs & inventory counters",
            "Booking DB: Reservation ledger, stepper status & timestamp audit logs",
            "Performance Indexing: B-Tree indexes on listing_id, customer_id & status"
        ]
    },
    {
        "num": 16,
        "title": "Frontend Architecture & UX Design",
        "subtitle": "Svelte 4 Reactive Store & Custom CSS Token System",
        "bullets": [
            "UI Framework: Svelte 4 compiler for zero-virtual-DOM performance",
            "State Management: Reactive stores (authStore.js, bookingStore.js)",
            "Visual Aesthetic: Modern dark theme, glassmorphism & gold accents (#d4af37)",
            "Component Library: Modular components (BookingModal, SpinWheelModal, Navbar)",
            "Accessibility (A11y): ARIA roles, high-contrast text & modal keyboard focus"
        ]
    },
    {
        "num": 17,
        "title": "API Specifications & Endpoints",
        "subtitle": "OpenAPI 3.0 Standardized Service Contracts",
        "bullets": [
            "GET /api/listings - Fetch 28 curated Ethiopian marketplace business listings",
            "POST /api/bookings?customerId=... - Execute atomic reservation transaction",
            "POST /api/wallet/topup - Reactive wallet credit deposit",
            "POST /api/spin/wheel - Execute promotional spin & register unique code",
            "Swagger UI: Live API documentation interactive sandbox at /swagger-ui.html"
        ]
    },
    {
        "num": 18,
        "title": "Quality Assurance & Test Suite",
        "subtitle": "Testing Pyramid & Defect Metric Analysis",
        "bullets": [
            "Unit Testing: JUnit 5 & Mockito testing core domain business rules",
            "Contract Testing: Pact framework verifying microservice API provider/consumer contracts",
            "Integration Testing: Testcontainers PostgreSQL & Kafka containerized integration tests",
            "End-to-End Testing: Playwright automated headless browser test flows",
            "Defect Zero Guarantee: Built-in fail-safe fallbacks for offline demo resiliency"
        ]
    },
    {
        "num": 19,
        "title": "Deployment & CI/CD Pipeline",
        "subtitle": "Docker Containerization & Cloud Hosting Strategy",
        "bullets": [
            "Containerization: Multi-stage Dockerfiles producing distroless Alpine images",
            "Orchestration: Kubernetes manifests (Deployment, Service, Ingress, HPA)",
            "Continuous Integration: GitHub Actions workflow building and testing images",
            "Cloud Deployment: Railway / Vercel containerized deployment",
            "Observability: Prometheus metrics scraper & Grafana dashboard visualizer"
        ]
    },
    {
        "num": 20,
        "title": "Conclusion & Project Roadmap",
        "subtitle": "Summary of Accomplishments & Future Extensions",
        "bullets": [
            "Accomplishment: Fully functional, hardened Ethiopian Service Marketplace platform",
            "Wallet Integrity: Reactive, fail-safe wallet payments and spin wheel charges",
            "Promo Engine: Unique single-use promo code generation and dynamic discount parsing",
            "Booking Flow: Seamless check-in/out date pickers and night multipliers",
            "Next Phase: Real-world Spring Boot backend REST integration & Telebirr Production API"
        ]
    }
]

# Write Slides Markdown document
slides_md_content = "# EasyService 20-Slide Presentation Deck\n\n"
for s in slides_data:
    slides_md_content += f"## Slide {s['num']}: {s['title']}\n"
    slides_md_content += f"**{s['subtitle']}**\n\n"
    for b in s['bullets']:
        slides_md_content += f"- {b}\n"
    slides_md_content += "\n---\n\n"

with open(SLIDES_MD_PATH, "w", encoding="utf-8") as f:
    f.write(slides_md_content)

print(f"[OK] Generated Markdown Presentation at: {SLIDES_MD_PATH}")

# Create PPTX Presentation
prs = Presentation()

# Set 16:9 Widescreen dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions
NAVY_DARK = RGBColor(15, 23, 42)
GOLD_ACCENT = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)
GRAY_TEXT = RGBColor(100, 116, 139)
DARK_TEXT = RGBColor(30, 41, 59)

blank_slide_layout = prs.slide_layouts[6]

for s_info in slides_data:
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 1. Background Shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.line.color.rgb = NAVY_DARK if s_info['num'] == 1 else WHITE
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_DARK if s_info['num'] == 1 else RGBColor(248, 250, 252)

    if s_info['num'] == 1:
        # Title Slide Layout (Dark Navy background with Gold typography)
        
        # Gold Accent Banner
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.5), Inches(11.333), Inches(0.1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = GOLD_ACCENT
        accent_bar.line.fill.background()

        # Title Box
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_info['title']
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = GOLD_ACCENT

        # Subtitle Box
        txBox_sub = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.333), Inches(0.8))
        tf_sub = txBox_sub.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = s_info['subtitle']
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = WHITE

        # Bullets / Details Box
        txBox_b = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.8))
        tf_b = txBox_b.text_frame
        tf_b.word_wrap = True
        for b_text in s_info['bullets']:
            p_b = tf_b.add_paragraph()
            p_b.text = f"*  {b_text}"
            p_b.font.size = Pt(15)
            p_b.font.color.rgb = RGBColor(226, 232, 240)
            p_b.space_after = Pt(8)

    else:
        # Standard Content Slide Layout
        
        # Top Header Banner Card
        header_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Inches(12.133), Inches(1.2))
        header_card.fill.solid()
        header_card.fill.fore_color.rgb = NAVY_DARK
        header_card.line.fill.background()

        # Slide Number Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.75), Inches(1.2), Inches(0.7))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GOLD_ACCENT
        badge.line.fill.background()
        tf_badge = badge.text_frame
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = f"Slide {s_info['num']}"
        p_badge.alignment = PP_ALIGN.CENTER
        p_badge.font.size = Pt(13)
        p_badge.font.bold = True
        p_badge.font.color.rgb = NAVY_DARK

        # Header Title Text
        txBox_t = slide.shapes.add_textbox(Inches(2.3), Inches(0.6), Inches(10.0), Inches(0.6))
        tf_t = txBox_t.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = s_info['title']
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = GOLD_ACCENT

        # Header Subtitle Text
        txBox_s = slide.shapes.add_textbox(Inches(2.3), Inches(1.1), Inches(10.0), Inches(0.5))
        tf_s = txBox_s.text_frame
        tf_s.word_wrap = True
        p_s = tf_s.paragraphs[0]
        p_s.text = s_info['subtitle']
        p_s.font.size = Pt(13)
        p_s.font.color.rgb = WHITE

        # Content Card Container
        content_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.9), Inches(12.133), Inches(5.1))
        content_card.fill.solid()
        content_card.fill.fore_color.rgb = WHITE
        content_card.line.color.rgb = RGBColor(226, 232, 240)

        # Bullets Box
        txBox_content = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.7))
        tf_c = txBox_content.text_frame
        tf_c.word_wrap = True
        
        for idx, b_text in enumerate(s_info['bullets']):
            p_c = tf_c.add_paragraph() if idx > 0 else tf_c.paragraphs[0]
            p_c.text = f"*   {b_text}"
            p_c.font.size = Pt(16)
            p_c.font.bold = True if idx == 0 else False
            p_c.font.color.rgb = DARK_TEXT
            p_c.space_after = Pt(16)

prs.save(SLIDES_PPTX_PATH)
print(f"[OK] Generated 20-Slide PowerPoint Presentation at: {SLIDES_PPTX_PATH}")
print("[OK] All document and presentation generation completed successfully!")
